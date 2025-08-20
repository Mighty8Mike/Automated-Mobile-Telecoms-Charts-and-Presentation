import os
import re
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
from matplotlib.image import imread
import matplotlib.image as mpimg
import matplotlib.ticker as mtick, matplotlib.ticker as ticker
from matplotlib.ticker import FuncFormatter
import seaborn as sns
import math
import threading
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Pt, Cm
from pptx.enum.shapes import MSO_SHAPE_TYPE  # ✅ Needed for shape type check
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.xmlchemy import OxmlElement
import tkinter as tk
from tkinter import filedialog
from tkinter.messagebox import showinfo
from io import BytesIO
from PIL import Image
import comtypes.client  # for converting slides to images (Windows only)
import shutil
from datetime import datetime




# === Paths ===
BASE_PATH = os.path.dirname(os.path.abspath(__file__))
CHARTS_PATH = os.path.join(BASE_PATH, 'Charts')
os.makedirs(CHARTS_PATH, exist_ok=True)
SLIDES_PATH = os.path.join(BASE_PATH, 'Slides')
os.makedirs(SLIDES_PATH, exist_ok=True)  
PRESENTATIONS_PATH = os.path.join(BASE_PATH, 'Presentations')
os.makedirs(PRESENTATIONS_PATH, exist_ok=True)

from Create_Charts import (select_image_files, display_images, read_entry)




# === Load & Clean Excel Data ===
def load_and_prepare_data(filename):
    df = pd.read_excel(filename, header=0)
    expected_columns = ['Key Indicator', 'WB Income Group', 'ITU Region', 'Country', 'Year', 'Value']
    missing = [col for col in expected_columns if col not in df.columns]
    if missing:
        raise ValueError(f"Missing columns in Excel: {missing}")

    df = df.dropna(subset=['Value'])
    df['Value'] = df['Value'].astype(str).str.replace('%', '').str.replace(',', '', regex=False)
    df['Value'] = pd.to_numeric(df['Value'], errors='coerce')
    df = df.dropna(subset=['Value'])

    def format_value(row):
        indicator = str(row['Key Indicator'])
        value = row['Value']
        if 'Penetration' in indicator:
            return f"{value:.1f}%"
        else:
            return f"{int(value):,}" if pd.notnull(value) else ""

    df['Key Indicator'] = df['Key Indicator'].astype(str)
    df['Formatted Value'] = df.apply(format_value, axis=1)
    return df



# Global chart parameters, to copy, not to cut 
plt.rcParams.update({
    'font.size': 20,
    'axes.titlesize': 22,
    'axes.labelsize': 20,
    'xtick.labelsize': 18,
    'ytick.labelsize': 18,
    'legend.fontsize': 18,
    'legend.title_fontsize': 20,
    'font.family': 'Calibri'  # or 'Arial', 'DejaVu Sans', etc.
})






    




# Compilation of content on the key slide layouts
SLIDE_WIDTH = Cm(33.867)
SLIDE_HEIGHT = Cm(19.05)
CRIMSON = RGBColor(192, 0, 0)



def add_textbox(slide, text, left, top, width, height, font_size,
                bold=False, italic=False, align="left", bullet=False,
                line_spacing=None, space_before=None, space_after=None,  font_color=RGBColor(0, 0, 0)):
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()

    # zero margins
    frame.margin_top = 0
    frame.margin_bottom = 0
    frame.margin_left = 0
    frame.margin_right = 0

    p = frame.paragraphs[0]
    p.text = text
    p.font.name = 'Calibri'
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.italic = italic

    if align == "center":
        p.alignment = PP_ALIGN.CENTER
    elif align == "right":
        p.alignment = PP_ALIGN.RIGHT
    else:
        p.alignment = PP_ALIGN.LEFT

    # *** Here's the bullet fix ***
    if bullet:
        p.bullet = True
        p.level = 0
        # The bullet character:
        # This sets a proper bullet char and font for PowerPoint
        run = p.runs[0]
        run.text = text  # Ensure run has the correct text
        run.font.name = 'Calibri'
        run.font.color.rgb = RGBColor(0, 0, 0)  # Black bullet text color
        # No need to use _element.set for bullet char here

    else:
        run = p.runs[0]
        run.font.color.rgb = RGBColor(0, 0, 0)  # Default black text

    return box



def add_line(slide, top, width, left, thickness):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(1))
    shape.fill.background()
    shape.line.color.rgb = CRIMSON
    shape.line.width = Pt(thickness)


def prepare_slides(slides_path=SLIDES_PATH, charts_path=CHARTS_PATH):
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    root = tk.Tk()
    root.withdraw()

    layout_choice = tk.simpledialog.askinteger(
        "Slide Layout",
        "Choose slide layout:\n1. Cover Slide\n2. Executive Summary\n3. 2-chart Slide\n4. 3-chart Slide",
        minvalue=1,
        maxvalue=4
    )

    if not layout_choice:
        return

    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank slide

    if layout_choice == 1:
        # Cover Slide
        add_textbox(slide, "Adjust Title", left=Cm(2), top=Cm(7), width=Cm(30), height=Cm(1.5), font_size=32, bold=True, align="center")
        add_line(slide, top=Cm(8.5), width=Cm(30), left=Cm(2), thickness=1.2)
        add_textbox(slide, "Adjust Text", left=Cm(2), top=Cm(17.5), width=Cm(30), height=Cm(1.5), font_size=16, bold=True, align="center")

    elif layout_choice == 2:
        # Executive Summary Slide
        add_textbox(slide, "Adjust Title", left=Cm(2), top=Cm(1), width=Cm(30), height=Cm(1.2), font_size=24, bold=True)
        add_line(slide, top=Cm(2.2), width=Cm(30), left=Cm(2), thickness=1)

        # Bullet list text box
        add_textbox(slide, text="Key Insight 1\nKey Insight 2\nKey Insight 3", left=Cm(2), top=Cm(2.7), width=Cm(30), height=Cm(12), font_size=14, bullet=True, line_spacing=Pt(12), space_before=Pt(0), space_after=Pt(0))

        

    elif layout_choice == 3:
        # 2-chart Slide
        # 1. Title
        add_textbox(slide, "Adjust Title", left=Cm(2), top=Cm(1), width=Cm(30), height=Cm(1.2), font_size=24, bold=True)
        title_box = add_textbox(slide, "Adjust Title", left=Cm(2), top=Cm(1), width=Cm(30), height=Cm(1.2), font_size=24, bold=True)
        title_right_x = title_box.left + title_box.width

        # 2. Title line directly under title
        add_line(slide, top=Cm(2.2), width=Cm(30), left=Cm(2), thickness=1)

        # 3. Key Message textbox 0.4 cm below the line
        add_textbox(slide, "Adjust Text", left=Cm(2), top=Cm(2.6), width=Cm(30), height=Cm(7), font_size=14, bold=True, bullet=True)

        

         # Ask for 2 chart images
        chart_files = filedialog.askopenfilenames(
            title="Select 2 charts",
            initialdir=CHARTS_PATH,
            filetypes=[("Image Files", "*.jpg *.jpeg *.png *.pdf")]
        )

        if len(chart_files) < 2:
            print("❗ Please select at least 2 charts.")
            return

        # === Left and Right Charts ===
        # Chart 1 (left-hand chart)
        top_y = Cm(2.6) + Cm(7) + Cm(0.5)   # Key Message from top + its height + 0.2 cm below Key Message box

        add_textbox(slide, "Adjust Header", left=Cm(2), top=top_y, width=Cm(14), height=Cm(0.6), font_size=12, bold=True)
        add_line(slide, top=top_y + Cm(0.7), width=Cm(14), left=Cm(2), thickness=1)
        slide.shapes.add_picture(chart_files[0], Cm(2), top_y + Cm(0.9), width=Cm(14), height=Cm(6))

        # Chart 2 (right-hand chart)
        # Align right edge of quadrant 2 to title's right edge
        chart_width = Cm(14)
        right_x = title_right_x - chart_width

        # Now apply this right_x:
        add_textbox(slide, "Adjust Header", left=right_x, top=top_y, width=Cm(14), height=Cm(0.6), font_size=12, bold=True)
        add_line(slide, top=top_y + Cm(0.7), width=Cm(14), left=right_x, thickness=1)
        slide.shapes.add_picture(chart_files[1], right_x, top_y + Cm(0.9), width=Cm(14), height=Cm(6))


        # Source text
        source_width = Cm(12)
        source_left = title_right_x - source_width  # align right edge

        add_textbox(slide, "Source: ITU", left=source_left, top=SLIDE_HEIGHT - Cm(1.5), width=source_width, height=Cm(1), font_size=8, italic=True, bold=True, align="right")


    elif layout_choice == 4:
        # 4-Quadrant Slide

        # 1. Title
        add_textbox(slide, "Adjust Title", left=Cm(2), top=Cm(1), width=Cm(30), height=Cm(1.2), font_size=24, bold=True)
        title_box = add_textbox(slide, "Adjust Title", left=Cm(2), top=Cm(1), width=Cm(30), height=Cm(1.2), font_size=24, bold=True)
        title_right_x = title_box.left + title_box.width

        # 2. Title line directly under title
        add_line(slide, top=Cm(2.2), width=Cm(30), left=Cm(2), thickness=1)

        # 3. Key Message textbox 0.4 cm below the line
        add_textbox(slide, "Adjust Text", left=Cm(2), top=Cm(2.6), width=Cm(30), height=Cm(0.8), font_size=14, bold=True, bullet=True)

        
        # Ask for 3 chart images
        chart_files = filedialog.askopenfilenames(
            title="Select 3 charts",
            initialdir=CHARTS_PATH,
            filetypes=[("Image Files", "*.jpg *.jpeg *.png *.pdf")]
        )

        if len(chart_files) < 3:
            print("❗ Please select at least 3 charts.")
            return

        # === TOP Quadrants ===
        # 4. Quadrant 1 (top-left text box)
        top_y = Cm(2.6) + Cm(0.8) + Cm(0.2)   # 0.2 cm below Key Message box

        add_textbox(slide, "Adjust Header", left=Cm(2), top=top_y, width=Cm(14), height=Cm(0.6), font_size=12, bold=True)
        add_line(slide, top=top_y + Cm(0.7), width=Cm(14), left=Cm(2), thickness=1)
        add_textbox(slide, "Adjust Text", left=Cm(2), top=top_y + Cm(0.9), width=Cm(14), height=Cm(4), font_size=11, bullet=True)

        # 5. Quadrant 2 (top-right chart)
        # Align right edge of quadrant 2 to title's right edge
        chart_width = Cm(14)
        right_x = title_right_x - chart_width

        # Now apply this right_x:
        add_textbox(slide, "Adjust Header", left=right_x, top=top_y, width=Cm(14), height=Cm(0.6), font_size=12, bold=True)
        add_line(slide, top=top_y + Cm(0.7), width=Cm(14), left=right_x, thickness=1)
        slide.shapes.add_picture(chart_files[0], right_x, top_y + Cm(0.9), width=Cm(14), height=Cm(6))


        # === BOTTOM Quadrants ===
        bottom_y = top_y + Cm(7.0)  # Ensure no overlap

        # 6. Quadrant 3 (bottom-left chart)
        add_textbox(slide, "Adjust Header", left=Cm(2), top=bottom_y, width=Cm(14), height=Cm(0.6), font_size=12, bold=True)
        add_line(slide, top=bottom_y + Cm(0.7), width=Cm(14), left=Cm(2), thickness=1)
        slide.shapes.add_picture(chart_files[1], Cm(2), bottom_y + Cm(0.9), width=Cm(14), height=Cm(6))

        # 7. Quadrant 4 (bottom-right chart)
        # Align right edge of quadrant 4 to title's right edge
        chart_width = Cm(14)
        right_x = title_right_x - chart_width

        # Now apply this right_x:
        add_textbox(slide, "Adjust Header", left=right_x, top=bottom_y, width=Cm(14), height=Cm(0.6), font_size=12, bold=True)
        add_line(slide, top=bottom_y + Cm(0.7), width=Cm(14), left=right_x, thickness=1)
        slide.shapes.add_picture(chart_files[2], right_x, bottom_y + Cm(0.9), width=Cm(14), height=Cm(6))


        # 8. Source text
        source_width = Cm(12)
        source_left = title_right_x - source_width  # align right edge

        add_textbox(slide, "Source: ITU", left=source_left, top=SLIDE_HEIGHT - Cm(1.5), width=source_width, height=Cm(1), font_size=8, italic=True, bold=True, align="right")





    # --- Save the slides ---
    # Ensure the Slides folder exists
    os.makedirs(SLIDES_PATH, exist_ok=True)

    # --- Create safe layout prefix ---
    layout_prefix = f"slide_layout_{layout_choice}"
    safe_prefix = re.sub(r'\W+', '_', layout_prefix)

    # --- Count existing files with the same prefix ---
    # Match pattern: slide_layout_3_1.pptx, slide_layout_3_2.pptx, etc.
    pattern = re.compile(rf'{re.escape(safe_prefix)}_(\d+)\.pptx')

    existing_indices = [
        int(match.group(1))
        for f in os.listdir(slides_path)
        if (match := pattern.match(f))
    ]

    next_index = max(existing_indices, default=0) + 1
    filename = f"{safe_prefix}_{next_index}.pptx"
    save_path = os.path.join(SLIDES_PATH, filename)

    # --- Save the slide status ---
    prs.save(save_path)
    showinfo("Slide Saved", f"✅ Slide saved to:\n{save_path}")

















def add_slide_number(slide, number, final_ppt):
    left = final_ppt.slide_width - Cm(2.5)
    top = final_ppt.slide_height - Cm(1)
    txBox = slide.shapes.add_textbox(left, top, Cm(2), Cm(1))
    tf = txBox.text_frame
    tf.text = f"{number}"
    p = tf.paragraphs[0]
    p.font.size = Pt(10)
    p.font.bold = True
    p.alignment = 2  # Right


def print_slides():
    print("\nSelect an option:")
    print("1. Create a new presentation from selected slides.")
    print("2. Insert slide(s) into an existing presentation.")
    option = input("Enter option number (1 or 2): ").strip()

    if option == "2":
        # --- Step 1: Choose the target presentation ---
        presentations = sorted(f for f in os.listdir(PRESENTATIONS_PATH) if f.endswith('.pptx'))
        if not presentations:
            print("❌ No presentations found in the Presentations folder.")
            return

        print("\nAvailable Presentations:")
        for i, f in enumerate(presentations, 1):
            print(f"{i}. {f}")

        try:
            pres_index = int(input("Select a presentation number to insert into: ")) - 1
            pres_path = os.path.join(PRESENTATIONS_PATH, presentations[pres_index])
        except (ValueError, IndexError):
            print("❌ Invalid selection.")
            return

        target_ppt = Presentation(pres_path)
        max_slide_number = len(target_ppt.slides)

        # --- Step 2: Choose slides to insert ---
        slide_files = sorted(f for f in os.listdir(SLIDES_PATH) if f.endswith('.pptx'))
        print("\nAvailable Slide Files:")
        for i, f in enumerate(slide_files, 1):
            print(f"{i}. {f}")

        user_input = input("Enter slide numbers to insert (e.g., 1,3-4): ")
        selected_slide_indices = parse_slide_input(user_input, len(slide_files))

        # --- Step 3: Select position(s) to insert after ---
        print(f"Presentation has {max_slide_number} slides.")

        # Show slide titles
        for i, slide in enumerate(target_ppt.slides, 1):
            title = None
            for shape in slide.shapes:
                if shape.has_text_frame and shape.text_frame.text.strip():
                    title = shape.text_frame.text.strip().replace("\n", " ")
                    break
            print(f"{i}. {title if title else '[No title]'}")

        insert_input = input(f"Insert after which slide number(s)? (e.g., 2,4): ")

        insert_after_indices = parse_slide_input(insert_input, max_slide_number)

        # Expand insert positions to match selected slides
        if len(insert_after_indices) == 1:
            insert_after_indices *= len(selected_slide_indices)
        elif len(insert_after_indices) != len(selected_slide_indices):
            print("❌ Number of insert positions must be either 1 or match the number of slides.")
            return

        # Process each (slide to insert, position) pair
        offset = 0
        for slide_index, insert_after in zip(selected_slide_indices, insert_after_indices):
            slide_path = os.path.join(SLIDES_PATH, slide_files[slide_index - 1])
            source_ppt = Presentation(slide_path)
            src_slide = source_ppt.slides[0]  # Always insert first slide only

            blank_slide_layout = target_ppt.slide_layouts[6]
            new_slide = target_ppt.slides.add_slide(blank_slide_layout)

            image_shapes = []
            for shape in src_slide.shapes:
                try:
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image_shapes.append({
                            "blob": shape.image.blob,
                            "left": shape.left,
                            "top": shape.top,
                            "width": shape.width,
                            "height": shape.height
                        })
                    else:
                        el = shape.element
                        new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                except Exception as e:
                    print(f"⚠️ Error copying shape: {e}")

            for img in image_shapes:
                new_slide.shapes.add_picture(
                    BytesIO(img["blob"]),
                    img["left"],
                    img["top"],
                    width=img["width"],
                    height=img["height"]
                )

            # Reorder slides: move inserted to correct index
            xml_slides = target_ppt.slides._sldIdLst
            slides = list(xml_slides)
            xml_slides.remove(slides[-1])
            insert_position = insert_after + offset
            xml_slides.insert(insert_position, slides[-1])
            offset += 1  # Account for growing slide list

        # Re-number slides (skip cover slides)
        # First: Remove old slide number textboxes
        for slide in target_ppt.slides:
            for shape in list(slide.shapes):  # convert to list to avoid iteration issues while deleting
                if shape.has_text_frame and shape.text_frame.text.strip().isdigit():
                    if shape.width == Cm(2) and shape.height == Cm(1):
                        slide.shapes._spTree.remove(shape._element)

        # Then: Re-number slides (skip covers)
        for i, slide in enumerate(target_ppt.slides, 1):
            is_cover = False
            for shape in slide.shapes:
                if shape.has_text_frame and "slide_layout_1_" in shape.text_frame.text:
                    is_cover = True
                    break
            if not is_cover:
                add_slide_number(slide, i, target_ppt)


        # Save updated presentation
        new_name = f"Updated_{presentations[pres_index]}"
        save_path = os.path.join(PRESENTATIONS_PATH, new_name)
        target_ppt.save(save_path)
        print(f"\n✅ Slides inserted and saved as {new_name} in 'Presentations' folder.")

    elif option == "1":
        # --- Step 1: Load available slide files ---
        files = sorted(f for f in os.listdir(SLIDES_PATH) if f.endswith('.pptx'))
        if not files:
            print("No slide files found in the Slides folder.")
            return

        print("\nAvailable Slide Files:")
        for i, f in enumerate(files, 1):
            print(f"{i}. {f}")

        # --- Step 2: User selects slides ---
        user_input = input("\nEnter slide numbers in desired order (e.g., 1,2,4-5): ")
        selected_indices = parse_slide_input(user_input, len(files))

        # --- Step 3: Compile new presentation ---
        final_ppt = Presentation()
        final_ppt.slide_width = Inches(13.33)
        final_ppt.slide_height = Inches(7.5)

        slide_counter = 1
        for index in selected_indices:
            path = os.path.join(SLIDES_PATH, files[index - 1])
            src_ppt = Presentation(path)

            for slide in src_ppt.slides:
                blank_slide_layout = final_ppt.slide_layouts[6]
                new_slide = final_ppt.slides.add_slide(blank_slide_layout)

                image_shapes = []
                for shape in slide.shapes:
                    try:
                        if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            image_shapes.append({
                                "blob": shape.image.blob,
                                "left": shape.left,
                                "top": shape.top,
                                "width": shape.width,
                                "height": shape.height
                            })
                        else:
                            el = shape.element
                            new_slide.shapes._spTree.insert_element_before(el, 'p:extLst')
                    except Exception as e:
                        print(f"⚠️ Error copying shape: {e}")

                for img in image_shapes:
                    new_slide.shapes.add_picture(
                        BytesIO(img["blob"]),
                        img["left"],
                        img["top"],
                        width=img["width"],
                        height=img["height"]
                    )

        # --- Step 4: Clean up old slide numbers ---
        for slide in final_ppt.slides:
            for shape in list(slide.shapes):
                if shape.has_text_frame and shape.text_frame.text.strip().isdigit():
                    if shape.width == Cm(2) and shape.height == Cm(1):
                        slide.shapes._spTree.remove(shape._element)

        # --- Step 5: Add slide numbers (skip cover slides) ---
        for i, slide in enumerate(final_ppt.slides, 1):
            is_cover = False
            for shape in slide.shapes:
                if shape.has_text_frame and "slide_layout_1_" in shape.text_frame.text:
                    is_cover = True
                    break
            if not is_cover:
                add_slide_number(slide, i, final_ppt)

        # --- Step 6: Save new presentation ---
        base_name = "Presentation"
        existing_files = [f for f in os.listdir(PRESENTATIONS_PATH) if f.startswith(base_name) and f.endswith('.pptx')]
        numbers = [int(f[len(base_name):-5]) for f in existing_files if f[len(base_name):-5].isdigit()]
        next_number = max(numbers + [0]) + 1
        final_name = f"{base_name}{next_number}.pptx"
        final_path = os.path.join(PRESENTATIONS_PATH, final_name)

        final_ppt.save(final_path)
        print(f"\n✅ Presentation saved as {final_name} in 'Presentations' folder.")


# Supporting helper function to parse slide input 
def parse_slide_input(input_str, max_num):
    slides = set()
    parts = input_str.replace(' ', '').split(',')
    for part in parts:
        if '-' in part:
            start, end = map(int, part.split('-'))
            slides.update(range(start, end + 1))
        else:
            slides.add(int(part))
    return sorted([s for s in slides if 1 <= s <= max_num])














def delete(): 
    folder_map = {
        "1": ("Charts", CHARTS_PATH),
        "2": ("Slides", SLIDES_PATH),
        "3": ("Presentations", PRESENTATIONS_PATH)
    }

    print("\nWhich folder would you like to manage?")
    for k, v in folder_map.items():
        print(f"{k}. {v[0]} ({v[1]})")
    choice = input("Enter your choice: ").strip()

    if choice not in folder_map:
        print("❌ Invalid folder selection.")
        return

    folder_name, folder_path = folder_map[choice]
    files = [f for f in os.listdir(folder_path) if os.path.isfile(os.path.join(folder_path, f))]
    if not files:
        print(f"No files found in {folder_name}.")
        return

    print(f"\nFiles in {folder_name}:")
    for idx, f in enumerate(files, 1):
        print(f"{idx}. {f}")

    if choice != "3":
        # For Charts and Slides — basic delete
        delete_input = input("Enter file numbers to delete (e.g., 1,3-4): ").strip()
        delete_indices = parse_selection_input(delete_input, len(files))
        for idx in delete_indices:
            file_to_delete = os.path.join(folder_path, files[idx - 1])
            os.remove(file_to_delete)
            print(f"✅ Deleted: {files[idx - 1]}")

    else:
        # === Presentation Menu ===
        print("\nPresentation management options:")
        print("1. Delete Slides from Presentation")
        print("2. Delete Entire Presentation")
        sub_choice = input("Enter your choice (1 or 2): ").strip()

        if sub_choice not in {"1", "2"}:
            print("❌ Invalid option.")
            return

        pres_num = input("Select presentation number: ").strip()
        if not pres_num.isdigit() or int(pres_num) < 1 or int(pres_num) > len(files):
            print("❌ Invalid presentation number.")
            return

        pres_file = files[int(pres_num) - 1]
        pres_path = os.path.join(folder_path, pres_file)

        if sub_choice == "1":
            # === Delete Slides ===
            prs = Presentation(pres_path)
            print(f"\nSlides in {pres_file}:")
            for i, slide in enumerate(prs.slides, 1):
                print(f"{i}. {slide.shapes.title.text if slide.shapes.title else 'Untitled Slide'}")

            del_input = input("Enter slide numbers to delete (e.g., 1,3-4): ").strip()
            indices_to_delete = sorted(parse_selection_input(del_input, len(prs.slides)), reverse=True)

            for idx in indices_to_delete:
                delete_slide(prs, idx - 1)
                print(f"✅ Deleted slide {idx}")

            prs.save(pres_path)
            print("✅ Presentation updated.")

        elif sub_choice == "2":
            # === Delete Entire Presentation ===
            os.remove(pres_path)
            print(f"✅ Deleted entire presentation: {pres_file}")


# === Helper functions for Delete, Presentations ===
def parse_selection_input(selection, max_val):
    indices = set()
    parts = selection.split(',')
    for part in parts:
        if '-' in part:
            start, end = part.split('-')
            indices.update(range(int(start), int(end) + 1))
        else:
            indices.add(int(part))
    return sorted([i for i in indices if 1 <= i <= max_val])

def delete_slide(prs, slide_index):
    slide_id = prs.slides._sldIdLst[slide_index].rId
    prs.part.drop_rel(slide_id)
    del prs.slides._sldIdLst[slide_index]








# === Module Guard ===
if __name__ == "__main__":
    print("This is a helper module. Please run ITU_Main.py instead.")
